from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .course_outline import CourseOutline
from .presentation_models import PresentationBuildReport, PresentationTheme
from .slide_models import SlideDeck

class PowerPointBuilder:
    def __init__(self, theme: PresentationTheme | None = None):
        self.theme = theme or PresentationTheme()

    def build(
        self,
        *,
        outline: CourseOutline,
        deck: SlideDeck,
        output_path: str | Path,
        artifact_root: str | Path = ".",
    ) -> PresentationBuildReport:
        output_path = Path(output_path)
        artifact_root = Path(artifact_root)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        if self.theme.aspect_ratio == "16:9":
            prs.slide_width = Inches(13.333333)
            prs.slide_height = Inches(7.5)

        figure_count = 0
        code_block_count = 0
        missing: list[str] = []
        warnings: list[str] = []

        self._add_title_slide(prs, outline)

        for slide_spec in deck.slides:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            self._set_title(slide, slide_spec.title)

            content_left = Inches(0.7)
            content_top = Inches(1.45)
            content_width = Inches(6.0)
            content_height = Inches(5.2)

            if slide_spec.figure_artifact:
                figure_path = artifact_root / slide_spec.figure_artifact
                if figure_path.exists():
                    slide.shapes.add_picture(
                        str(figure_path),
                        Inches(7.0),
                        Inches(1.55),
                        width=Inches(5.6),
                    )
                    figure_count += 1
                    content_width = Inches(5.9)
                else:
                    missing.append(slide_spec.figure_artifact)

            text_box = slide.shapes.add_textbox(
                content_left,
                content_top,
                content_width,
                content_height,
            )
            frame = text_box.text_frame
            frame.word_wrap = True
            frame.clear()

            for index, bullet in enumerate(slide_spec.bullets):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
                paragraph.font.name = self.theme.body_font
                paragraph.font.size = Pt(self.theme.body_font_size_pt)
                paragraph.space_after = Pt(8)

            if slide_spec.code_artifact:
                code_path = artifact_root / slide_spec.code_artifact
                if code_path.exists():
                    code = code_path.read_text(encoding="utf-8")
                    self._add_code_block(slide, code)
                    code_block_count += 1
                else:
                    missing.append(slide_spec.code_artifact)

            if slide_spec.references:
                self._add_reference_footer(slide, slide_spec.references)

            if slide_spec.speaker_notes:
                notes = slide.notes_slide.notes_text_frame
                notes.text = slide_spec.speaker_notes

        self._add_references_slide(prs, outline.references)
        prs.save(output_path)

        if len(deck.slides) == 0:
            warnings.append("Slide deck contains no teaching slides.")

        return PresentationBuildReport(
            output_path=str(output_path),
            slide_count=len(prs.slides),
            figure_count=figure_count,
            code_block_count=code_block_count,
            missing_artifacts=tuple(sorted(set(missing))),
            warnings=tuple(warnings),
        )

    def _add_title_slide(self, prs: Presentation, outline: CourseOutline) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = outline.title
        slide.placeholders[1].text = (
            f"Audience: {outline.audience}\n"
            f"Duration: {outline.total_duration_minutes} minutes"
        )
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = self.theme.title_font
            paragraph.font.size = Pt(self.theme.title_font_size_pt + 6)

    def _set_title(self, slide, title: str) -> None:
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.8)
        )
        paragraph = title_box.text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.name = self.theme.title_font
        paragraph.font.size = Pt(self.theme.title_font_size_pt)
        paragraph.font.bold = True

    def _add_code_block(self, slide, code: str) -> None:
        box = slide.shapes.add_textbox(
            Inches(0.7), Inches(4.65), Inches(5.9), Inches(2.0)
        )
        frame = box.text_frame
        frame.word_wrap = False
        frame.margin_left = Pt(6)
        frame.margin_right = Pt(6)
        paragraph = frame.paragraphs[0]
        paragraph.text = code[:2500]
        paragraph.font.name = self.theme.code_font
        paragraph.font.size = Pt(self.theme.code_font_size_pt)

    def _add_reference_footer(self, slide, references: tuple[str, ...]) -> None:
        box = slide.shapes.add_textbox(
            Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.3)
        )
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = "Sources: " + ", ".join(references)
        paragraph.font.name = self.theme.body_font
        paragraph.font.size = Pt(8)
        paragraph.alignment = PP_ALIGN.RIGHT

    def _add_references_slide(
        self,
        prs: Presentation,
        references: tuple[str, ...],
    ) -> None:
        if not references:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        self._set_title(slide, "References")
        box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5)
        )
        frame = box.text_frame
        frame.word_wrap = True

        for index, reference in enumerate(references):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = reference
            paragraph.font.name = self.theme.body_font
            paragraph.font.size = Pt(16)
