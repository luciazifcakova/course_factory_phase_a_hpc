from pathlib import Path

from pptx import Presentation

from course_factory import (
    CourseModule,
    CourseOutline,
    JobContext,
    Lesson,
    PowerPointBuilder,
    PowerPointBuilderAgent,
    Slide,
    SlideDeck,
)

def build_outline():
    return CourseOutline(
        title="Introduction to ggplot2",
        audience="Beginners",
        language="English",
        modules=(
            CourseModule(
                module_id="basics",
                title="Basics",
                lessons=(
                    Lesson(
                        lesson_id="scatter",
                        title="Scatter plots",
                        duration_minutes=30,
                        practical=True,
                        required_packages=("ggplot2",),
                    ),
                ),
            ),
        ),
        learning_objectives=("Create plots",),
        required_packages=("ggplot2",),
        total_duration_minutes=60,
        references=("DOC-1",),
    )

def build_deck():
    return SlideDeck(
        course_title="Introduction to ggplot2",
        slides=(
            Slide(
                slide_id="scatter-1",
                lesson_id="scatter",
                title="Scatter plots",
                bullets=("Use geom_point()", "Map variables in aes()"),
                speaker_notes="Demonstrate with iris.",
                references=("DOC-1",),
                code_artifact="scripts/scatter.R",
                figure_artifact="figures/scatter.png",
            ),
        ),
    )

def test_builder_creates_powerpoint(tmp_path):
    scripts = tmp_path / "scripts"
    scripts.mkdir()
    (scripts / "scatter.R").write_text(
        "library(ggplot2)\nggplot(iris, aes(Sepal.Length, Sepal.Width)) + geom_point()",
        encoding="utf-8",
    )

    output = tmp_path / "course.pptx"
    report = PowerPointBuilder().build(
        outline=build_outline(),
        deck=build_deck(),
        output_path=output,
        artifact_root=tmp_path,
    )

    assert output.exists()
    assert report.slide_count == 3
    assert report.code_block_count == 1
    assert report.figure_count == 0
    assert report.missing_artifacts == ("figures/scatter.png",)

    prs = Presentation(output)
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text == "Introduction to ggplot2"

def test_powerpoint_builder_agent(tmp_path):
    outline = build_outline()
    deck = build_deck()
    context = JobContext.create(user_request="Build deck").model_copy(
        update={
            "state": {
                "course_outline": outline.model_dump(mode="json"),
                "slide_deck": deck.model_dump(mode="json"),
            }
        }
    )

    result = PowerPointBuilderAgent(
        output_dir=tmp_path / "presentations",
        artifact_root=tmp_path,
    ).run(context)

    assert result.status.value == "success"
    assert Path(result.outputs["presentation_path"]).exists()
    assert result.metrics["presentation_slides"] == 3
